import os
from datetime import time, datetime
from docx import Document
from pptx import Presentation
import olefile
from trash.comp_image.image_parser import ImageProcessor
from hwp5.hwp5txt import Hwp5File

class DocumentParser:
    def __init__(self, processor: ImageProcessor):
        self.processor = processor

    def parse_docx(self, file_path):
        """
        DOCX 파일을 파싱합니다.
        - 텍스트는 각 문단별로 추출하여 리스트에 저장
        - 이미지 추출은 ImageProcessor.extract_images_docx()
        메타데이터(파일명, 파일 크기, 최종 수정 시간)도 함께 반환
        """
        try:
            doc = Document(file_path)
        except Exception as e:
            raise Exception(f"Error loading DOCX file: {e}")

        contents = []
        current_problem = None
        image_counter = 1

        # 이미지 추출
        image_map = self.processor.extract_images_docx(doc)

        for element in self._iter_docx_elements(doc):
            if isinstance(element, Paragraph):
                text = element.text.strip()
                if not text:
                    continue


                if text.startswith("문제"):


    def parse_ppt(self, file_path):
        """
        PPTX 파일을 파싱합니다.
        - 텍스트는 각 슬라이드별로 추출하여 리스트에 저장
        - 이미지 추출은 ImageProcessor.extract_images_pptx()
        메타데이터(파일명, 파일 크기, 최종 수정 시간)도 함께 반환
        """
        try:
            prs = Presentation(file_path)
        except Exception as e:
            raise Exception(f"Error loading PPTX file: {e}")

        contents = []
        image_count = self.processor.extract_images_pptx(prs)
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        slide_text.append(text)
            if slide_text:
                contents.append({"type": "slide", "content": "\n".join(slide_text)})
        return {
            "metadata": self._get_metadata(file_path),
            "contents": contents,
            "images": image_count
        }

    def parse_hwp(self, file_path):
        import re
        # HWP 텍스트 추출: olefile 방식으로 PrvText 스트림 사용
        text_content = self.extract_hwp_text(file_path)  # 아래에 정의된 함수 사용
        paragraphs = [p.strip() for p in text_content.split('\x0D') if p.strip()]

        # 이미지 추출은 ImageProcessor에서 처리한 결과 사용
        image_infos = self.processor.extract_images_hwp(file_path)

        contents = []
        current_problem = None
        problem_pattern = re.compile(r'(문제\s*\d+|\[Q\d+\])')
        image_ref_pattern = re.compile(r'\(그림\s*\d+\)')

        for para in paragraphs:
            if problem_pattern.match(para):
                if current_problem:
                    contents.append(current_problem)
                current_problem = {
                    "type": "problem",
                    "number": problem_pattern.search(para).group(),
                    "content": [para],
                    "images": []  # 추후 이미지 연결
                }
            else:
                # 만약 이미지 참조가 있다면 매칭
                image_refs = image_ref_pattern.findall(para)
                for ref in image_refs:
                    img_num = int(re.search(r'\d+', ref).group())
                    if 0 < img_num <= len(image_infos):
                        current_problem['images'].append(image_infos[img_num - 1]['path'])
                if current_problem:
                    current_problem['content'].append(para)
                else:
                    contents.append({"type": "text", "content": para})
        if current_problem:
            contents.append(current_problem)

        return {
            "metadata": self._get_metadata(file_path),
            "contents": contents,
            "images": image_infos
        }

    def _iter_hwp_paragraphs(self, hwp):
        text_content = []
        try:
            for section in hwp.bodytext.sections():
                for para in section.paragraphs():
                    yield para.text.replace('\n', ' ').strip()
        except Exception as e:
            raise RuntimeError(f"HWP 문단 추출 오류: {str(e)}")


    def extract_hwp_text(self, file_path):
        try:
            with olefile.OleFileIO(file_path) as ole:
                if not ole.exists('PrvText'):
                    raise RuntimeError("HWP 파일에 'PrvText' 스트림이 없습니다.")
                encoded_text = ole.openstream('PrvText').read()
                return encoded_text.decode('utf-16le')
        except Exception as e:
            raise RuntimeError(f"HWP 텍스트 추출 오류: {str(e)}")


    # 메타데이터 생성
    def _get_metadata(self, file_path):
        stat = os.stat(file_path)
        return {
            "file_name": os.path.basename(file_path),
            "file_size": stat.st_size,
            "modified_time": datetime.fromtimestamp(stat.st_mtime).isoformat()
        }
