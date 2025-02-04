import os

import olefile
from PIL import Image
from pytesseract import pytesseract

pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

# ---------------------------
# Image Processing Class
# ---------------------------

class ImageProcessor:

    def __init__(self, output_dir="output_images"):
        self.output_dir = output_dir
        os.makedirs(self.output_dir, exist_ok=True)

    def process_extracted_images(self):
        image_texts = []
        for img_file in os.listdir(self.output_dir):
            if img_file.lower().endswith((".png", ".jpg", ".jpeg")):
                text = self.ocr_process(os.path.join(self.output_dir, img_file))
                if text:
                    image_texts.append({
                        "image": img_file,
                        "text": text
                    })
        return image_texts

    def extract_images_docx(self, doc):
        """
        :param doc: images?
        :return: save image blob data into file
        """
        image_count = 0
        try:
            # doc.part.rels : Docx file's relationship recursive
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    image_count += 1
                    image_data = rel.target_part.blob
                    img_path = os.path.join(self.output_dir, f"docx_image_{image_count}.png")
                    with open(img_path, "wb") as f:
                        f.write(image_data)
        except Exception as e:
            print(f"Error extracting images from DOCX: {e}")
        return image_count

    def extract_images_pptx(self, prs):
        """
        PPTX 파일에서 이미지를 추출합니다.
        PPTX의 각 슬라이드 내의 picture shape을 찾아 blob 데이터를 저장하고,
        저장된 이미지의 수를 반환합니다.
        """
        image_count = 0
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        try:
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image_count += 1
                        img_path = os.path.join(self.output_dir, f"pptx_image_{image_count}.png")
                        with open(img_path, "wb") as f:
                            f.write(shape.image.blob)
        except Exception as e:
            print(f"Error extracting images from PPTX: {e}")
        return image_count

    def extract_images_hwp(self, hwp_file_path):
        """HWP 파일에서 이미지 추출 (olefile.OleFileIO 사용)"""
        from hwp5.xmlmodel import Hwp5File
        hwp = Hwp5File(hwp_file_path)
        image_infos = []

        try:
            ole = olefile.OleFileIO(hwp_file_path)
            if not ole.exists('BinData'):
                print("HWP 파일에 'BinData' 스트림이 없습니다. 이미지 추출 건너뜀.")
                ole.close()
                return image_infos  # 빈 리스트 반환

            # ole.listdir()로 스트림 목록 순회 (entry는 예: ['BinData', 'BIN0'])
            for entry in ole.listdir(streams=True, storages=False):
                if entry[0] == 'BinData':
                    try:
                        data = ole.openstream(entry).read()
                        # 이미지 파일 서명 확인 (JPEG, PNG)
                        if data.startswith(b'\xff\xd8'):
                            ext = 'jpg'
                        elif data.startswith(b'\x89PNG\r\n\x1a\n'):
                            ext = 'png'
                        else:
                            continue  # 이미지가 아니면 건너뜀

                        img_index = len(image_infos) + 1
                        img_path = os.path.join(self.output_dir, f"hwp_image_{img_index}.{ext}")
                        with open(img_path, "wb") as f:
                            f.write(data)

                        # OCR 처리 후 결과 포함
                        ocr_text = self.ocr_process(img_path)
                        image_infos.append({
                            "id": img_index,
                            "path": img_path,
                            "position": None,  # 위치 정보가 있다면 여기 채워넣기
                            "ocr_text": ocr_text
                        })
                    except Exception as inner_e:
                        print(f"스트림 {entry} 처리 중 오류: {inner_e}")
            ole.close()
            return image_infos
        except Exception as e:
            raise RuntimeError(f"HWP 이미지 처리 오류: {str(e)}")

    def ocr_process(self, image_path):
            """
            주어진 이미지 파일에 대해 OCR 수행하여 텍스트를 추출
            """
            try:
                image = Image.open(image_path)
                return pytesseract.image_to_string(image, lang='kor+eng').strip()
            except Exception as e:
                print(f"OCR Error for {image_path}: {e}")
                return ""

