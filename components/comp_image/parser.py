import os
from datetime import time
from docx import Document
from pptx import Presentation
import olefile
from hwp5 import hwp5txt
from components.comp_image.extractor import ImageProcessor


class DocumentParser:
    def __init__(self, processor: ImageProcessor):
        self.processor = processor

    def parse_docx(self, file_path):
        """
        DOCX 파일을 파싱합니다.
        - 텍스트는 각 문단별로 추출하여 리스트에 저장
        - 이미지 추출은 ImageProcessor.extract_images_docx()
        메타데이터(파일명, 파일 크기, 최종 수정 시간)도 함께 반환
        """
        try:
            doc = Document(file_path)
        except Exception as e:
            raise Exception(f"Error loading DOCX file: {e}")


        contents = []
        image_count = self.processor.extract_images_docx(doc)
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                contents.append({"type": "text", "content": text})
        return {
            "metadata": self._get_metadata(file_path),
            "contents": contents,
            "images": image_count
        }


    def parse_ppt(self, file_path):
        """
        PPTX 파일을 파싱합니다.
        - 텍스트는 각 슬라이드별로 추출하여 리스트에 저장
        - 이미지 추출은 ImageProcessor.extract_images_pptx()
        메타데이터(파일명, 파일 크기, 최종 수정 시간)도 함께 반환
        """
        try:
            prs = Presentation(file_path)
        except Exception as e:
            raise Exception(f"Error loading PPTX file: {e}")

        contents = []
        image_count = self.processor.extract_images_pptx(prs)
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        slide_text.append(text)
            if slide_text:
                contents.append({"type": "slide", "content": "\n".join(slide_text)})
        return {
            "metadata": self._get_metadata(file_path),
            "contents": contents,
            "images": image_count
        }


    def parse_hwp(self, file_path):
        """
        HWP 파일을 파싱합니다.
        - 텍스트는 각 문단별로 추출하여 리스트에 저장
        - 이미지 추출은 ImageProcessor.extract_images_hwp()
        메타데이터(파일명, 파일 크기, 최종 수정 시간)도 함께 반환
        """
        try:
            contents = []
            image_count = self.processor.extract_images_hwp(file_path)

            # 텍스트 추출
            with open(file_path, "rb") as f:
                hwp_text = hwp5txt(f)
                text_content = hwp_text.get_text()

                # 문단 단위 분리
                paragraphs = [p.strip() for p in text_content.split('\x0D') if p.strip()]
                for para in paragraphs:
                    contents.append({"type": "text", "content": para})

            return {
                "metadata": self._get_metadata(file_path),
                "contents": contents,
                "images": image_count
            }
        except Exception as e:
            raise Exception(f"HWP 파싱 오류: {e}")


    def _get_metadata(self, file_path):
        """
        파일의 메타데이터를 반환합니다.
        """
        file_name = os.path.basename(file_path)
        file_size = os.path.getsize(file_path)
        modified_time = time.ctime(os.path.getmtime(file_path))
        return {
            "file_name": file_name,
            "file_size": file_size,
            "modified_time": modified_time
        }