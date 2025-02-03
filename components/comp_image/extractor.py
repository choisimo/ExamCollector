import os

import olefile
from PIL.Image import Image
from pytesseract import pytesseract


# ---------------------------
# Image Processing Class
# ---------------------------

class ImageProcessor:

    def __init__(self, output_dir="output_images"):
        self.output_dir = output_dir
        os.makedirs(self.output_dir, exist_ok=True)

    def process_extracted_images(self):
        image_texts = []
        for img_file in os.listdir(self.output_dir):
            if img_file.lower().endswith((".png", ".jpg", ".jpeg")):
                text = self.ocr_process(os.path.join(self.output_dir, img_file))
                if text:
                    image_texts.append({
                        "image": img_file,
                        "text": text
                    })
        return image_texts

    def extract_images_docx(self, doc):
        """
        :param doc: images?
        :return: save image blob data into file
        """
        image_count = 0
        try:
            # doc.part.rels : Docx file's relationship recursive
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    image_count += 1
                    image_data = rel.target_part.blob
                    img_path = os.path.join(self.output_dir, f"docx_image_{image_count}.png")
                    with open(img_path, "wb") as f:
                        f.write(image_data)
        except Exception as e:
            print(f"Error extracting images from DOCX: {e}")
        return image_count

    def extract_images_pptx(self, prs):
        """
        PPTX 파일에서 이미지를 추출합니다.
        PPTX의 각 슬라이드 내의 picture shape을 찾아 blob 데이터를 저장하고,
        저장된 이미지의 수를 반환합니다.
        """
        image_count = 0
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        try:
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image_count += 1
                        img_path = os.path.join(self.output_dir, f"pptx_image_{image_count}.png")
                        with open(img_path, "wb") as f:
                            f.write(shape.image.blob)
        except Exception as e:
            print(f"Error extracting images from PPTX: {e}")
        return image_count

    def extract_images_hwp(self, hwp_file_path):
        """
        HWP 파일에서 이미지를 추출합니다.
        :param hwp_file_path:
        :return:
        """
        image_count = 0
        try:
            with olefile.OleFileIO(hwp_file_path) as hwp:
                if 'Bindata' not in hwp.listdir():
                    return 0

            for entry in hwp.listdir():
                if entry[0].startswith('BinData/BIN'):
                    bindata = hwp.openstream(entry).read()
                    if bindata[:4] == b'\xFF\xD8\xFF\xE0':  # JPEG 시그니처
                        ext = 'jpg'
                    elif bindata[:8] == b'\x89\x50\x4E\x47\x0D\x0A\x1A\x0A':  # PNG 시그니처
                        ext = 'png'
                    else:
                        continue

                    image_count += 1
                    img_path = os.path.join(
                        self.output_dir,
                        f"hwp_image_{image_count}.{ext}"
                    )
                    with open(img_path, "wb") as f:
                        f.write(bindata)
            return image_count
        except Exception as e:
            print(f"HWP 이미지 추출 오류: {e}")
            return 0


def ocr_process(self, image_path):
        """
        주어진 이미지 파일에 대해 OCR 수행하여 텍스트를 추출
        """
        try:
            image = Image.open(image_path)
            text = pytesseract.image_to_string(image, lang='kor+eng')
            return text.strip()
        except Exception as e:
            print(f"OCR Error for {image_path}: {e}")
            return ""

